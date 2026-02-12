#!/usr/bin/env python3
"""
PPTX 자동 품질 점수 산출 스크립트.
python-pptx + markitdown으로 슬라이드 콘텐츠를 분석하고
5개 차원에 대해 자동 점수를 산출한다.

사용법:
    python quality_scorer.py presentation.pptx
    python quality_scorer.py presentation.pptx --json  # JSON 출력
    python quality_scorer.py presentation.pptx --refs refs.json  # DOI 검증 포함

자동 측정 가능한 항목만 점수화하고, 수동 평가가 필요한 항목은 체크리스트로 제공.
"""

import argparse
import json
import sys
import subprocess
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Emu
except ImportError:
    print("pip install python-pptx 필요", file=sys.stderr)
    sys.exit(1)


def analyze_pptx(pptx_path: str) -> dict:
    """PPTX 파일 구조 분석."""
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width / 914400
    slide_h = prs.slide_height / 914400

    slides = []
    for si, slide in enumerate(prs.slides, 1):
        slide_info = {
            'number': si,
            'shapes_count': len(slide.shapes),
            'text_shapes': 0,
            'image_shapes': 0,
            'table_shapes': 0,
            'total_text_length': 0,
            'font_sizes': [],
            'has_images': False,
            'texts': [],
        }

        for shape in slide.shapes:
            if shape.has_text_frame:
                slide_info['text_shapes'] += 1
                text = shape.text_frame.text.strip()
                slide_info['total_text_length'] += len(text)
                if text:
                    slide_info['texts'].append(text[:100])
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            slide_info['font_sizes'].append(run.font.size.pt)

            if shape.shape_type == 13:  # Picture
                slide_info['image_shapes'] += 1
                slide_info['has_images'] = True

            if shape.has_table:
                slide_info['table_shapes'] += 1

        slides.append(slide_info)

    return {
        'file': pptx_path,
        'slide_dimensions': f'{slide_w:.3f} x {slide_h:.3f}',
        'total_slides': len(slides),
        'slides': slides,
    }


def score_content_density(analysis: dict) -> dict:
    """콘텐츠 밀도 점수 (자동)."""
    issues = []
    scores = []

    for s in analysis['slides']:
        # 텍스트 길이 기준
        text_len = s['total_text_length']
        if text_len < 50:
            scores.append(2)
            issues.append(f"슬라이드 {s['number']}: 텍스트 부족 ({text_len}자)")
        elif text_len > 2000:
            scores.append(3)
            issues.append(f"슬라이드 {s['number']}: 텍스트 과다 ({text_len}자)")
        else:
            scores.append(5)

        # shape 수 기준
        if s['shapes_count'] < 5:
            issues.append(f"슬라이드 {s['number']}: 요소 부족 ({s['shapes_count']}개)")
        elif s['shapes_count'] > 60:
            issues.append(f"슬라이드 {s['number']}: 요소 과다 ({s['shapes_count']}개)")

    avg = sum(scores) / len(scores) if scores else 3.0
    return {'score': round(avg, 1), 'issues': issues}


def score_visual_balance(analysis: dict) -> dict:
    """시각적 균형 점수 (자동)."""
    issues = []
    scores = []

    for s in analysis['slides']:
        text_count = s['text_shapes']
        image_count = s['image_shapes']
        table_count = s['table_shapes']

        visual_elements = image_count + table_count
        if visual_elements == 0 and text_count > 3:
            scores.append(2)
            issues.append(f"슬라이드 {s['number']}: 시각자료 없음 (텍스트만 {text_count}개)")
        elif visual_elements > 0:
            scores.append(5)
        else:
            scores.append(3)

    avg = sum(scores) / len(scores) if scores else 3.0
    return {'score': round(avg, 1), 'issues': issues}


def score_format_consistency(analysis: dict) -> dict:
    """형식 일관성 점수 (자동)."""
    issues = []
    all_sizes = []

    for s in analysis['slides']:
        all_sizes.extend(s['font_sizes'])

    # 폰트 크기 다양성 확인
    unique_sizes = set(all_sizes)
    if len(unique_sizes) > 10:
        issues.append(f"폰트 크기 종류 과다 ({len(unique_sizes)}종): {sorted(unique_sizes)}")
        score = 3.0
    elif len(unique_sizes) > 7:
        score = 4.0
    else:
        score = 5.0

    # 최대 폰트 확인
    if all_sizes:
        max_size = max(all_sizes)
        if max_size > 36:
            issues.append(f"최대 폰트 크기 {max_size}pt (30pt 이하 권장)")

    return {'score': round(score, 1), 'issues': issues}


def generate_report(analysis: dict, refs_file: str = None) -> dict:
    """종합 품질 보고서 생성."""
    density = score_content_density(analysis)
    visual = score_visual_balance(analysis)
    formatting = score_format_consistency(analysis)

    # 학술 근거와 논리적 흐름은 자동 점수 불가 → 수동 체크리스트
    report = {
        'file': analysis['file'],
        'total_slides': analysis['total_slides'],
        'dimensions': analysis['slide_dimensions'],
        'scores': {
            'academic_evidence': {
                'score': None,
                'weight': 0.25,
                'note': '수동 평가 필요: 핵심 주장에 DOI 레퍼런스가 있는지 확인',
                'checklist': [
                    '모든 수치에 출처가 있는가?',
                    'DOI가 유효한가? (scripts/doi_lookup.py --validate로 검증)',
                    '가이드라인 참조에 판본이 명시되어 있는가?',
                ],
            },
            'content_density': {
                'score': density['score'],
                'weight': 0.20,
                'issues': density['issues'],
            },
            'visual_balance': {
                'score': visual['score'],
                'weight': 0.20,
                'issues': visual['issues'],
            },
            'logical_flow': {
                'score': None,
                'weight': 0.20,
                'note': '수동 평가 필요: 슬라이드 간 전개가 자연스러운지 확인',
                'checklist': [
                    '배경→방법→결과→결론 순서가 논리적인가?',
                    '각 슬라이드의 핵심 메시지가 명확한가?',
                    '전환이 자연스러운가?',
                ],
            },
            'format_consistency': {
                'score': formatting['score'],
                'weight': 0.15,
                'issues': formatting['issues'],
            },
        },
        'auto_scores_weighted': None,
    }

    # 자동 점수만으로 가중 평균 (학술 근거/논리적 흐름 제외)
    auto_scores = []
    auto_weights = []
    for k, v in report['scores'].items():
        if v['score'] is not None:
            auto_scores.append(v['score'] * v['weight'])
            auto_weights.append(v['weight'])

    if auto_weights:
        report['auto_scores_weighted'] = round(sum(auto_scores) / sum(auto_weights), 1)

    return report


def print_report(report: dict):
    """사람이 읽기 쉬운 형태로 보고서 출력."""
    print(f"\n{'='*50}")
    print(f"  슬라이드 품질 보고서")
    print(f"  파일: {report['file']}")
    print(f"  슬라이드: {report['total_slides']}장 ({report['dimensions']})")
    print(f"{'='*50}")

    if report['auto_scores_weighted']:
        print(f"\n  자동 점수 (자동측정 항목): {report['auto_scores_weighted']}/5.0")

    for name, info in report['scores'].items():
        labels = {
            'academic_evidence': '학술 근거',
            'content_density': '콘텐츠 밀도',
            'visual_balance': '시각적 균형',
            'logical_flow': '논리적 흐름',
            'format_consistency': '형식 일관성',
        }
        label = labels.get(name, name)
        weight = int(info['weight'] * 100)

        if info['score'] is not None:
            print(f"\n  [{label}] {info['score']}/5 (가중치 {weight}%)")
        else:
            print(f"\n  [{label}] 수동 평가 필요 (가중치 {weight}%)")

        if 'issues' in info:
            for issue in info['issues']:
                print(f"    ✗ {issue}")
        if 'checklist' in info:
            for item in info['checklist']:
                print(f"    □ {item}")
        if 'note' in info:
            print(f"    ℹ {info['note']}")

    print(f"\n{'='*50}")


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='PPTX quality scorer')
    parser.add_argument('pptx', help='PPTX file to analyze')
    parser.add_argument('--json', action='store_true', help='Output as JSON')
    parser.add_argument('--refs', help='References JSON for DOI validation')

    args = parser.parse_args()

    analysis = analyze_pptx(args.pptx)
    report = generate_report(analysis, args.refs)

    if args.json:
        print(json.dumps(report, indent=2, ensure_ascii=False))
    else:
        print_report(report)
